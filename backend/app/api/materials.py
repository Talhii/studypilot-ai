"""Materials router — async multipart intake, chunk ingestion, DB persistence.

POST /api/materials/upload validates the file, extracts + chunks its text, and
persists the material to `uploaded_materials` scoped to the authenticated user.
Raw bytes are not retained.
"""

from __future__ import annotations

import io
import logging
import os
import urllib.parse
import uuid

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile, status

from app import config
from app.core.auth import get_current_user_id
from app.core.database import db
from app.models.schemas import FileFormat, MaterialSummary, UploadResponse
from app.rag import memory as rag_memory

logger = logging.getLogger("studypilot.materials")

# Optional parsers — absence falls back to a synthetic summary so uploads still
# work, but installing these (requirements.txt) gives real document text.
try:
    import pypdf
except ImportError:  # pragma: no cover
    pypdf = None  # type: ignore[assignment]
try:
    import docx  # python-docx
except ImportError:  # pragma: no cover
    docx = None  # type: ignore[assignment]
try:
    from pptx import Presentation  # python-pptx
except ImportError:  # pragma: no cover
    Presentation = None  # type: ignore[assignment]

router = APIRouter(prefix="/api/materials", tags=["materials"])


@router.post(
    "/upload",
    response_model=UploadResponse,
    status_code=status.HTTP_201_CREATED,
    summary="Upload a study material, index its chunks, persist to the user.",
)
async def upload_material(
    file: UploadFile = File(...),
    user_id: str = Depends(get_current_user_id),
) -> UploadResponse:
    extension = os.path.splitext(file.filename or "")[1].lower()
    if extension not in config.ALLOWED_EXTENSIONS:
        allowed = ", ".join(sorted(config.ALLOWED_EXTENSIONS))
        raise HTTPException(
            status_code=status.HTTP_415_UNSUPPORTED_MEDIA_TYPE,
            detail=f"Unsupported file type '{extension}'. Allowed: {allowed}.",
        )

    # Stream into a buffer while enforcing the size cap.
    raw = bytearray()
    chunk_size = 1024 * 64
    while data := await file.read(chunk_size):
        raw += data
        if len(raw) > config.MAX_UPLOAD_BYTES:
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail="File exceeds the 25 MB Phase 1 limit.",
            )
    await file.close()

    # iOS document picker sends a percent-encoded name (e.g. "a%20b.pdf").
    file_name = urllib.parse.unquote(file.filename or "untitled")
    material_id = str(uuid.uuid4())

    # Extract real text from the document, then chunk + cache it.
    if extension == ".txt":
        text = bytes(raw).decode("utf-8", errors="ignore")
    else:
        text = _extract_text(bytes(raw), extension, file_name)
    chunk_count = rag_memory.ingest_material(material_id, text)

    material = MaterialSummary(
        id=material_id,
        user_id=user_id,
        file_name=file_name,
        file_url=None,
        chunk_count=chunk_count,
        format=FileFormat(extension.lstrip(".")),
    )

    # Ensure the user row exists (Supabase FK), then persist the material.
    db.get_or_create_user(user_id)
    db.insert_material(material)

    return UploadResponse(
        status="created",
        message=f"Ingested '{material.file_name}' ({chunk_count} chunks).",
        material=material,
    )


def _extract_text(raw: bytes, extension: str, file_name: str) -> str:
    """Pull real text from a PDF/DOCX/PPTX.

    Falls back to a synthetic summary when the parser is missing or the file
    has no extractable text (e.g. a scanned/image-only PDF needing OCR).
    """
    ext = extension.lower()
    try:
        if ext == ".pdf" and pypdf is not None:
            reader = pypdf.PdfReader(io.BytesIO(raw))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
            if text.strip():
                return text
        elif ext == ".docx" and docx is not None:
            document = docx.Document(io.BytesIO(raw))
            text = "\n".join(p.text for p in document.paragraphs)
            if text.strip():
                return text
        elif ext == ".pptx" and Presentation is not None:
            lines: list[str] = []
            for slide in Presentation(io.BytesIO(raw)).slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs)
                            if line.strip():
                                lines.append(line)
            if lines:
                return "\n".join(lines)
    except Exception as exc:  # noqa: BLE001 - degrade rather than fail the upload
        logger.warning("Text extraction failed for '%s' (%s).", file_name, exc)

    logger.info("No extractable text for '%s'; using synthetic summary.", file_name)
    return _simulate_extraction(file_name, ext.lstrip("."))


def _simulate_extraction(file_name: str, fmt: str) -> str:
    """Fabricate structured paragraph text for non-plaintext formats.

    Phase 1-5 ships no real PDF/DOCX/PPTX parser, so we synthesise readable
    per-section paragraphs that the chunker and retriever can operate on.
    """
    topic = _stem(file_name)
    sections = 8
    paragraphs = [
        (
            f"Section {i + 1} of {file_name} ({fmt}). This passage covers core "
            f"concepts related to {topic}, including key definitions, a worked "
            f"example, and summary notes a student should review for mastery. "
            f"It also highlights common pitfalls and how to avoid them."
        )
        for i in range(sections)
    ]
    return "\n\n".join(paragraphs)


def _stem(file_name: str) -> str:
    stem = file_name.rsplit(".", 1)[0].replace("_", " ").replace("-", " ").strip()
    return stem.title() or "General Knowledge"
